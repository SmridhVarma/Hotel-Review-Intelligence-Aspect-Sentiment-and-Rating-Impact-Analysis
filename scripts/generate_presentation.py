"""
generate_presentation.py
========================
Generates the BT5153 Group 19 presentation as a styled .pptx file.

Design: deep navy (#1C3557) + warm sand-gold (#C4A882) on off-white (#FAF9F5).
Minimalist — no gradients, no shadows, clean whitespace.

Usage:
    pip install python-pptx
    python scripts/generate_presentation.py

Output:
    presentation/BT5153_Group19_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design tokens ──────────────────────────────────────────────────────────────
BG       = RGBColor(250, 249, 245)  # warm off-white background
NAVY     = RGBColor(28,  53,  87)   # deep navy — primary
GOLD     = RGBColor(196, 168, 130)  # warm sand-gold — accent
DARK_TXT = RGBColor(44,  44,  44)   # body text
MID_TXT  = RGBColor(107, 107, 107)  # secondary / muted text
WHITE    = RGBColor(255, 255, 255)
CARD_BG  = RGBColor(235, 242, 248)  # light blue-grey for card fills
NAVY_DK  = RGBColor(18,  33,  55)   # darker navy for title slide texture

W = Inches(13.33)
H = Inches(7.5)


# ── Primitive helpers ──────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, italic=False,
             color=DARK_TXT, align=PP_ALIGN.LEFT,
             font="Calibri Light"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def chrome(slide, section_label=""):
    """Left navy bar + top gold rule + optional section label. Applied to all content slides."""
    add_rect(slide, Inches(0),    Inches(0), Inches(0.07), H,           NAVY)
    add_rect(slide, Inches(0.07), Inches(0), W - Inches(0.07), Inches(0.04), GOLD)
    if section_label:
        add_text(slide, section_label.upper(),
                 Inches(0.25), Inches(0.1), Inches(8), Inches(0.35),
                 size=8, bold=True, color=GOLD, font="Calibri")


def slide_title(slide, title):
    add_text(slide, title,
             Inches(0.3), Inches(0.45), Inches(12.5), Inches(0.95),
             size=32, bold=True, color=NAVY)
    add_rect(slide, Inches(0.3), Inches(1.35), Inches(12.5), Inches(0.025), GOLD)


# ── Slide templates ────────────────────────────────────────────────────────────

def bullet_slide(prs, title, items, section_label=""):
    """
    Standard content slide.
    items: list of strings.
      Prefix "#" for a bold sub-header within the slide.
      Prefix "##" for a muted smaller point.
      Empty string "" for vertical spacing.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    chrome(slide, section_label)
    slide_title(slide, title)

    tb = slide.shapes.add_textbox(Inches(0.35), Inches(1.55), Inches(12.2), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after  = Pt(3)

        if item == "":
            p.add_run().text = " "
            p.add_run().font.size = Pt(6)
            continue

        if item.startswith("## "):
            r = p.add_run()
            r.text = item[3:]
            r.font.size  = Pt(15)
            r.font.color.rgb = MID_TXT
            r.font.italic = True
            r.font.name   = "Calibri"
        elif item.startswith("# "):
            r = p.add_run()
            r.text = item[2:]
            r.font.size  = Pt(17)
            r.font.bold  = True
            r.font.color.rgb = NAVY
            r.font.name  = "Calibri Light"
        else:
            r = p.add_run()
            r.text = item
            r.font.size  = Pt(17)
            r.font.color.rgb = DARK_TXT
            r.font.name  = "Calibri"

    return slide


def section_divider(prs, number, title, subtitle=""):
    """Full-bleed navy section break slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, Inches(0),  Inches(5.6), W, Inches(0.055), GOLD)
    add_rect(slide, Inches(0.7), Inches(1.3), Inches(0.055), Inches(3.0), GOLD)

    add_text(slide, f"{number:02d}",
             Inches(0.8), Inches(1.2), Inches(3), Inches(2.6),
             size=88, color=RGBColor(255, 255, 255), font="Calibri Light")
    add_text(slide, title,
             Inches(1.9), Inches(2.0), Inches(10.5), Inches(1.4),
             size=40, bold=True, color=WHITE, font="Calibri Light")
    if subtitle:
        add_text(slide, subtitle,
                 Inches(1.9), Inches(3.3), Inches(10.5), Inches(0.8),
                 size=19, color=GOLD, font="Calibri Light")
    return slide


def two_col_slide(prs, title, l_head, l_items, r_head, r_items, section_label=""):
    """Two equal card columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    chrome(slide, section_label)
    slide_title(slide, title)

    col_w = Inches(5.95)
    gap   = Inches(0.35)
    top   = Inches(1.55)
    ht    = Inches(5.55)
    l_x   = Inches(0.28)
    r_x   = l_x + col_w + gap

    for x, head, items in [(l_x, l_head, l_items), (r_x, r_head, r_items)]:
        add_rect(slide, x, top, col_w, ht, CARD_BG)
        add_rect(slide, x, top, col_w, Inches(0.055), NAVY)
        add_text(slide, head,
                 x + Inches(0.18), top + Inches(0.1), col_w - Inches(0.25), Inches(0.5),
                 size=15, bold=True, color=NAVY, font="Calibri Light")
        add_rect(slide, x + Inches(0.18), top + Inches(0.58), Inches(0.055), ht - Inches(0.75), GOLD)

        tb = slide.shapes.add_textbox(
            x + Inches(0.35), top + Inches(0.65),
            col_w - Inches(0.5), ht - Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(3)
            if item == "":
                p.add_run().text = " "
                continue
            r = p.add_run()
            r.text = item
            r.font.size  = Pt(15)
            r.font.color.rgb = DARK_TXT if not item.startswith('"') else MID_TXT
            r.font.name  = "Calibri"
            r.font.italic = item.startswith('"')

    return slide


def placeholder_slide(prs, title, instruction, caption="", section_label=""):
    """Slide with a large shaded placeholder box for a figure or screenshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    chrome(slide, section_label)
    slide_title(slide, title)

    add_rect(slide, Inches(1.4), Inches(1.55), Inches(10.5), Inches(5.15), CARD_BG)
    add_rect(slide, Inches(1.4), Inches(1.55), Inches(10.5), Inches(0.04), GOLD)

    add_text(slide, instruction,
             Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.8),
             size=13, italic=True, color=MID_TXT,
             align=PP_ALIGN.CENTER, font="Calibri")
    if caption:
        add_text(slide, caption,
                 Inches(1.4), Inches(6.8), Inches(10.5), Inches(0.35),
                 size=10, italic=True, color=MID_TXT,
                 align=PP_ALIGN.CENTER, font="Calibri")
    return slide


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── 01 Title ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_rect(s, Inches(0),            Inches(0),            Inches(0.08), H,            GOLD)
    add_rect(s, Inches(0),            H - Inches(0.08),     W,            Inches(0.08), GOLD)
    add_rect(s, W - Inches(3.6),      Inches(0),            Inches(3.6),  H,            NAVY_DK)
    add_rect(s, W - Inches(3.6),      Inches(0),            Inches(0.04), H,            RGBColor(50, 80, 115))

    add_text(s, "BT5153  Applied Machine Learning for Business Analytics  |  AY 2025/2026",
             Inches(0.28), Inches(0.35), Inches(12.5), Inches(0.38),
             size=10, color=GOLD, font="Calibri Light")
    add_rect(s, Inches(0.28), Inches(0.72), Inches(9.2), Inches(0.025), RGBColor(70, 100, 135))

    add_text(s, "Hotel Review Intelligence",
             Inches(0.28), Inches(1.05), Inches(9.8), Inches(1.3),
             size=50, bold=True, color=WHITE, font="Calibri Light")
    add_text(s, "Aspect Sentiment and Rating Impact Analysis",
             Inches(0.28), Inches(2.25), Inches(9.8), Inches(0.75),
             size=24, color=GOLD, font="Calibri Light")
    add_text(s, "A Multi-Module LangGraph Pipeline for Automated Hotel Review Analysis",
             Inches(0.28), Inches(2.95), Inches(9.8), Inches(0.55),
             size=15, color=RGBColor(175, 185, 200), font="Calibri Light")

    add_rect(s, Inches(0.28), Inches(3.6), Inches(4.8), Inches(0.025), RGBColor(80, 110, 145))
    add_text(s, "Group 19",
             Inches(0.28), Inches(3.75), Inches(4), Inches(0.4),
             size=14, bold=True, color=GOLD, font="Calibri")
    add_text(s,
             "Smridh Varma  ·  Siddarth Mahesh  ·  [Member 3]  ·  [Member 4]  ·  [Member 5]  ·  [Member 6]",
             Inches(0.28), Inches(4.15), Inches(9.5), Inches(0.38),
             size=12, color=RGBColor(190, 195, 205), font="Calibri Light")

    # NUS logo placeholder (right panel)
    add_rect(s, W - Inches(3.1), Inches(2.8), Inches(2.3), Inches(1.1), NAVY_DK)
    add_text(s, "[ NUS Logo ]",
             W - Inches(3.1), Inches(3.1), Inches(2.3), Inches(0.5),
             size=11, italic=True, color=RGBColor(90, 115, 150),
             align=PP_ALIGN.CENTER, font="Calibri Light")

    # ── 02 Agenda ────────────────────────────────────────────────────────────────
    bullet_slide(prs, "Agenda", [
        "01  Problem & Motivation",
        "02  System Overview",
        "03  Dataset",
        "04  Methodology: ABSA Pipeline",
        "05  Methodology: Rating Impact Modelling",
        "06  Methodology: Conversational Agent",
        "07  Evaluation",
        "08  Limitations & Future Work",
        "09  Conclusion & Demo",
    ], section_label="Overview")

    # ── 03 Section: Problem ───────────────────────────────────────────────────────
    section_divider(prs, 1, "Problem & Motivation",
                    "Why does unstructured review data remain underutilised?")

    # ── 04 Problem Statement ──────────────────────────────────────────────────────
    bullet_slide(prs, "The Problem", [
        "Hotels receive thousands of reviews — but lack tools to act on them systematically",
        "Reviews mix comments across multiple service aspects in a single entry",
        "Manual analysis is inefficient, inconsistent, and does not scale",
        "Managers see a low score but cannot identify which specific issues are responsible",
        "",
        "# The Business Question",
        "How can hotels automatically convert large volumes of textual feedback into structured insights and identify which service issues most affect ratings?",
        "",
        "## Secondary question: can a conversational interface make these insights accessible to non-technical operators?",
    ], section_label="Problem")

    # ── 05 Section: System Overview ───────────────────────────────────────────────
    section_divider(prs, 2, "System Overview",
                    "Three modules, one end-to-end pipeline")

    # ── 06 Pipeline Overview ─────────────────────────────────────────────────────
    placeholder_slide(
        prs,
        "What We Built",
        ("[ Pipeline diagram ]\n\n"
         "data.csv  →  Stage 1: Preprocess  →  Stage 2: Aspect Extraction\n"
         "→  Stage 3: Sentiment Assignment  →  Stage 4: Rating Impact Modelling\n"
         "→  Stage 5: ChromaDB Ingestion  →  LangGraph Agent  →  Gradio UI"),
        "End-to-end pipeline: raw reviews to conversational insights",
        section_label="System Overview"
    )

    # ── 07 Three Modules ──────────────────────────────────────────────────────────
    two_col_slide(
        prs, "Three Modules, One Goal",
        "ML Pipeline  (Modules A + B)",
        [
            "Module A — ABSA",
            "Extract aspect-level sentiment across 6 service dimensions from 515,738 reviews without manual annotation",
            "",
            "Module B — Rating Impact",
            "Model which aspects most drive guest ratings using regression + SHAP value decomposition",
            "",
            "Output feeds the agent as structured context",
        ],
        "Conversational Agent  (Module C)",
        [
            "Module C — RAG Agent",
            "LangGraph-based chatbot grounded in real review evidence and model outputs",
            "",
            "Hotel owners query insights in natural language",
            "",
            "Supports 4 query types:",
            "Evidence  ·  Prioritisation  ·  Mismatch  ·  Segment",
            "",
            "All answers cited from actual reviews",
        ],
        section_label="System Overview"
    )

    # ── 08 Section: Dataset ───────────────────────────────────────────────────────
    section_divider(prs, 3, "Dataset",
                    "515,738 reviews  ·  1,492 hotels  ·  6 European cities")

    # ── 09 Dataset ───────────────────────────────────────────────────────────────
    two_col_slide(
        prs, "European Hotel Reviews Dataset",
        "Key Statistics",
        [
            "515,738 reviews across 1,492 hotels",
            "Cities: Amsterdam, Barcelona, London, Milan, Paris, Vienna",
            "Review dates: Aug 2015 to Aug 2017",
            "227 unique reviewer nationalities",
            "",
            "Reviewer_Score: 0–10 continuous rating",
            "Mean ≈ 8.4  ·  Left-skewed distribution",
            "",
            "[ Add reviewer score histogram here ]",
        ],
        "Why This Dataset Works",
        [
            "Pre-separated fields:",
            "Positive_Review and Negative_Review",
            "",
            "This enables deterministic polarity labelling without any manual annotation",
            "",
            "Tags field encodes trip type:",
            "Couple  ·  Family  ·  Business  ·  Solo  ·  Group",
            "",
            "Enables reviewer-segment-aware retrieval in the agent",
        ],
        section_label="Dataset"
    )

    # ── 10 Section: Methodology ───────────────────────────────────────────────────
    section_divider(prs, 4, "Methodology",
                    "Stages 1–4: from raw text to structured features")

    # ── 11 Stages 1–3: ABSA ───────────────────────────────────────────────────────
    bullet_slide(prs, "Stages 1–3: Aspect-Based Sentiment Analysis", [
        "# Stage 1: Text Preparation",
        "Combine Positive_Review + Negative_Review  →  sentence tokenise  →  tag each sentence with source polarity and reviewer segment",
        "",
        "# Stage 2: Aspect Extraction via LDA",
        "LDA topic modelling discovers vocabulary co-occurrence clusters for 6 target aspects",
        "Aspects: Cleanliness  ·  Staff  ·  Location  ·  Noise  ·  Food  ·  Room",
        "Output: aspect keyword dictionary (topic-word distributions, manually curated)",
        "",
        "# Stage 3: Sentiment Assignment",
        "Keyword match sentence  →  aspect  ·  Sentiment from source field:  +1 positive,  -1 negative,  0 no match",
        "Outputs: aspect_sentences.csv (for agent)  +  review_features.csv (for modelling)",
    ], section_label="Methodology — ABSA")

    # ── 12 Stage 4: Rating Impact ─────────────────────────────────────────────────
    two_col_slide(
        prs, "Stage 4: Rating Impact Modelling",
        "Regression Models",
        [
            "Linear Regression",
            "Interpretable baseline — coefficients directly readable as aspect impact weights",
            "",
            "Tree-based Ensemble",
            "Captures nonlinear interactions between aspect features",
            "",
            "Evaluation: RMSE, MAE, R²",
            "80/20 train-test split",
            "",
            "[ Add RMSE / R² results here ]",
        ],
        "SHAP Value Decomposition",
        [
            "Theoretically grounded feature attribution (Shapley values)",
            "",
            "Global ranking",
            "Mean absolute SHAP across all 515k reviews",
            "",
            "Per-hotel ranking",
            "Run global model on each hotel's review subset",
            "",
            "Output: shap_summary.json",
            "__global__ entry + one entry per hotel",
            "",
            "[ Add SHAP bar chart + per-hotel heatmap here ]",
        ],
        section_label="Methodology — Rating Impact"
    )

    # ── 13 Section: Agent ─────────────────────────────────────────────────────────
    section_divider(prs, 5, "Conversational Agent",
                    "LangGraph  ·  ChromaDB  ·  HyDE  ·  GPT-4o")

    # ── 14 LangGraph DAG ──────────────────────────────────────────────────────────
    placeholder_slide(
        prs,
        "LangGraph DAG Architecture",
        ("[ LangGraph DAG diagram ]\n\n"
         "query_classifier  →  segment_filter  →  (conditional routing)\n\n"
         "Evidence / Segment path:     hyde_expander  →  evidence_retriever  →  context_merger\n"
         "Prioritisation / Mismatch:   summary_retriever  →  context_merger\n\n"
         "context_merger  →  response_generator  →  state_manager  →  END"),
        "8-node DAG with conditional routing by query type. State persists across conversation turns.",
        section_label="Agent"
    )

    # ── 15 Two ChromaDB Collections ───────────────────────────────────────────────
    two_col_slide(
        prs, "Hierarchical Retrieval: Two Collections",
        "evidence_store",
        [
            "~1 million sentence embeddings",
            "Model: text-embedding-3-small (OpenAI)",
            "",
            "Metadata per document:",
            "  hotel_name  ·  aspect  ·  sentiment",
            "  reviewer_segment  ·  reviewer_score",
            "",
            "Answers: why?",
            "Evidence queries + segment-specific queries",
            "Returns cited review excerpts",
        ],
        "summary_store",
        [
            "One document per hotel + __global__ entry",
            "Human-readable SHAP narrative per hotel",
            "",
            "Example document text:",
            '"For Hotel Arena: cleanliness has the largest negative impact (-0.42), staff the largest positive (+0.31)..."',
            "",
            "Answers: what to prioritise?",
            "Prioritisation queries",
            "Macro-level impact ranking",
        ],
        section_label="Agent"
    )

    # ── 16 HyDE + Segmentation ────────────────────────────────────────────────────
    two_col_slide(
        prs, "Smarter Retrieval: HyDE + Segmentation",
        "Hypothetical Document Embeddings",
        [
            "Problem: short user queries sit far from dense review text in embedding space",
            "",
            "HyDE solution:",
            "1. GPT-4o generates a 2–3 sentence hypothetical review that would answer the query",
            "2. Embed the hypothetical — not the raw query",
            "3. Use that embedding for ChromaDB similarity search",
            "",
            "Better semantic match against real reviews",
            "No fine-tuning or labelled data required",
        ],
        "Reviewer Segmentation",
        [
            "Tags field normalised to 5 segments:",
            "Business  ·  Couple  ·  Family  ·  Solo  ·  Group",
            "",
            "segment_filter node detects segment mention in query",
            "Passes ChromaDB where-filter to evidence_retriever",
            "",
            "Enables targeted queries:",
            '"What do families complain about?"',
            '"How do business travelers rate our staff?"',
            "",
            "Falls back to unfiltered retrieval if no segment detected",
        ],
        section_label="Agent"
    )

    # ── 17 Groundedness ───────────────────────────────────────────────────────────
    bullet_slide(prs, "Groundedness Constraint", [
        "Every factual claim in a response must be traceable to a retrieved chunk or summary document",
        "",
        "# Confidence Mechanism",
        "context_merger computes mean cosine similarity across retrieved chunks",
        "If mean similarity < 0.5, or both retrieval sources return nothing: low_confidence = True",
        "response_generator switches to fallback prompt: cannot answer confidently — please rephrase",
        "",
        "# Citations",
        "Source sentence  ·  hotel name  ·  aspect  ·  reviewer segment  ·  similarity score",
        "Surfaced alongside every response in the Gradio UI",
        "",
        "# Supported Query Types",
        "Evidence  ·  Prioritisation  ·  Mismatch detection  ·  Reviewer segment queries",
    ], section_label="Agent")

    # ── 18 Section: Evaluation ────────────────────────────────────────────────────
    section_divider(prs, 6, "Evaluation",
                    "Three modules, three sets of metrics")

    # ── 19 Evaluation ─────────────────────────────────────────────────────────────
    bullet_slide(prs, "Evaluation Results", [
        "# Module A — Aspect Classification",
        "Metrics: Accuracy, Precision, Recall, F1 per aspect + macro average",
        "Validation: stratified random sample of [N] manually labelled sentences",
        "[ Add per-aspect classification table + aspect frequency chart here ]",
        "",
        "# Module B — Rating Impact Modelling",
        "Metrics: RMSE, MAE, R² on held-out test set + SHAP rank stability across time slices",
        "[ Add regression results table + SHAP bar chart + per-hotel heatmap here ]",
        "",
        "# Module C — Conversational Agent",
        "Metrics: Groundedness  ·  Relevance (GPT-4o judge)  ·  Context accuracy  ·  Resolution rate",
        "Evaluated on [N] manually curated queries across all 4 query types",
        "[ Add agent evaluation scores table here ]",
    ], section_label="Evaluation")

    # ── 20 Demo ───────────────────────────────────────────────────────────────────
    placeholder_slide(
        prs,
        "Live Demo",
        ("[ Screenshot of Gradio UI ]\n\n"
         "Hotel dropdown: specific hotel selected\n"
         "Query: 'Which issue should we fix first to improve ratings?'\n"
         "Response: grounded answer citing SHAP summary + real review excerpts\n"
         "Citations panel visible below response"),
        "Hotel owner queries the agent; every claim is backed by retrieved evidence",
        section_label="Demo"
    )

    # ── 21 Section: Limitations ───────────────────────────────────────────────────
    section_divider(prs, 7, "Limitations & Future Work",
                    "Honest assessment and next steps")

    # ── 22 Limitations ────────────────────────────────────────────────────────────
    two_col_slide(
        prs, "Limitations & Future Work",
        "Current Limitations",
        [
            "6 fixed aspects only — pricing, parking, amenities excluded",
            "",
            "Keyword matching is brittle to synonyms, negations, and co-mentions",
            "",
            "Source-field polarity introduces noise — some reviewers mislabel fields",
            "",
            "Confidence threshold (0.5) set heuristically",
            "",
            "English-only — no multilingual support",
            "",
            "Spam and extreme ratings treated as aggregate signal only",
        ],
        "Future Work",
        [
            "Fine-tuned sequence labelling model for ABSA — handles negation and implicit aspects",
            "",
            "Nationality-based segmentation using Reviewer_Nationality",
            "",
            "Empirical calibration of confidence threshold against evaluation set",
            "",
            "Human-in-the-loop feedback for low-quality retrievals",
            "",
            "Multilingual extension via translation preprocessing",
        ],
        section_label="Limitations"
    )

    # ── 23 Conclusion ─────────────────────────────────────────────────────────────
    bullet_slide(prs, "Key Takeaways", [
        "515,738 unstructured reviews converted into structured aspect-level operational insights — without manual annotation",
        "",
        "SHAP-based impact ranking identifies which service aspects most drive guest ratings at global and per-hotel level",
        "",
        "LangGraph RAG agent with HyDE retrieval enables natural language querying of evidence-grounded insights",
        "",
        "Reviewer-segment filtering surfaces targeted insights for specific guest types",
        "",
        "# Directly answers the project brief:",
        "How can hotels automatically convert large volumes of customer feedback into structured insights and identify which service issues most affect ratings?",
    ], section_label="Conclusion")

    # ── 24 Thank You ──────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0),            Inches(0.08), H,            GOLD)
    add_rect(s, Inches(0), H - Inches(0.08),     W,            Inches(0.08), GOLD)

    add_text(s, "Thank You",
             Inches(1.5), Inches(1.8), Inches(10), Inches(1.6),
             size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri Light")
    add_text(s, "Questions?",
             Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
             size=30, color=GOLD,
             align=PP_ALIGN.CENTER, font="Calibri Light")
    add_rect(s, Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.025), GOLD)
    add_text(s, "Group 19  |  BT5153  |  AY 2025/2026",
             Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
             size=13, color=RGBColor(155, 165, 180),
             align=PP_ALIGN.CENTER, font="Calibri Light")

    # ── Save ──────────────────────────────────────────────────────────────────────
    os.makedirs("presentation", exist_ok=True)
    out = os.path.join("presentation", "BT5153_Group19_Presentation.pptx")
    prs.save(out)
    print(f"Saved  ->  {out}")
    print(f"Slides ->  {len(prs.slides)}")


if __name__ == "__main__":
    build()
